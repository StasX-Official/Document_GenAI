"""Document and presentation templating utilities."""
from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import List, Optional, Sequence

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor

try:
    from pptx import Presentation
    from pptx.util import Inches as PPTInches

    PPTX_AVAILABLE = True
except Exception:  # pragma: no cover - optional dependency
    PPTX_AVAILABLE = False
    Presentation = None  # type: ignore


class DocumentTemplate:
    def apply_template(self, doc: Document, content: str, prompt: str, model: str) -> Document:
        self._reset_document(doc)
        header = doc.add_heading("AI Generated Document", level=1)
        header.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_heading("About", level=2)
        metadata = doc.add_paragraph()
        metadata.add_run("Generated on: ").bold = True
        metadata.add_run(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
        metadata.add_run("\nPrompt: ").bold = True
        prompt_run = metadata.add_run(prompt)
        prompt_run.italic = True
        metadata.add_run("\nModel: ").bold = True
        metadata.add_run(model)

        doc.add_paragraph("").paragraph_format.space_after = Pt(12)
        doc.add_heading("Generated Content", level=2)

        for paragraph_text in filter(str.strip, content.split("\n\n")):
            paragraph = doc.add_paragraph()
            paragraph.paragraph_format.space_after = Pt(10)
            paragraph.paragraph_format.first_line_indent = Inches(0.25)
            paragraph.add_run(paragraph_text)

        doc.add_paragraph("").paragraph_format.space_before = Pt(18)
        footer = doc.add_paragraph(f"Generated using {model.split(':')[0].upper()} AI")
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer_run = footer.runs[0]
        footer_run.font.size = Pt(9)
        footer_run.font.color.rgb = RGBColor(120, 120, 120)
        return doc

    def add_image(self, doc: Document, image_bytes: bytes, width_inches: float = 6.0, caption: Optional[str] = None) -> Document:
        image_stream = BytesIO(image_bytes)
        paragraph = doc.add_paragraph()
        run = paragraph.add_run()
        run.add_picture(image_stream, width=Inches(width_inches))
        if caption:
            caption_paragraph = doc.add_paragraph(caption)
            caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            caption_run = caption_paragraph.runs[0]
            caption_run.font.size = Pt(9)
            caption_run.italic = True
        return doc

    def add_table(self, doc: Document, data: Sequence[Sequence[str]], headers: Optional[Sequence[str]] = None) -> Document:
        if headers:
            table = doc.add_table(rows=1, cols=len(headers))
            header_cells = table.rows[0].cells
            for index, header in enumerate(headers):
                header_cells[index].text = str(header)
            for row in data:
                cells = table.add_row().cells
                for index, cell in enumerate(row):
                    cells[index].text = str(cell)
        else:
            if not data:
                return doc
            table = doc.add_table(rows=0, cols=len(data[0]))
            for row in data:
                cells = table.add_row().cells
                for index, cell in enumerate(row):
                    cells[index].text = str(cell)
        return doc

    def create_presentation(self, title: str, slides: List[dict]) -> Optional[bytes]:
        if not PPTX_AVAILABLE:
            return None
        presentation = Presentation()
        title_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        for slide_payload in slides:
            layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = slide_payload.get("title", "")
            if slide_payload.get("content"):
                slide.shapes.placeholders[1].text = slide_payload["content"]
            if slide_payload.get("image"):
                image_stream = BytesIO(slide_payload["image"])
                slide.shapes.add_picture(image_stream, PPTInches(1), PPTInches(1), width=PPTInches(8))
        output = BytesIO()
        presentation.save(output)
        return output.getvalue()

    def create_presentation_from_text(
        self,
        content: str,
        prompt: str,
        model: str,
        image_bytes: Optional[bytes] = None,
    ) -> Optional[bytes]:
        if not PPTX_AVAILABLE:
            return None

        slides: List[dict] = []
        segments = [segment.strip() for segment in content.split("\n\n") if segment.strip()]
        if not segments:
            return None

        slides.append({
            "title": "Executive Summary",
            "content": segments[0][:1500],
            "image": image_bytes,
        })

        for idx, segment in enumerate(segments[1:], start=1):
            title, body = self._split_segment(segment, idx)
            slides.append({"title": title, "content": body})

        slides.append({
            "title": "Prompt & Model",
            "content": f"Prompt: {prompt}\nModel: {model}",
        })

        return self.create_presentation("AI Generated Presentation", slides)

    def render_html(self, content: str, prompt: str, model: str, image_bytes: Optional[bytes]) -> str:
        image_tag = ""
        if image_bytes:
            import base64

            encoded = base64.b64encode(image_bytes).decode("utf-8")
            image_tag = f'<img src="data:image/png;base64,{encoded}" alt="Generated" class="hero" />'

        paragraphs = "".join(f"<p>{segment.strip()}</p>" for segment in content.split("\n\n") if segment.strip())
        return f"""
<!DOCTYPE html>
<html lang=\"en\">
<head>
  <meta charset=\"utf-8\" />
  <title>AI Generated Document</title>
  <style>
    body {{ font-family: 'Segoe UI', sans-serif; margin: 2rem auto; max-width: 920px; line-height: 1.6; color: #1f2933; }}
    h1, h2 {{ color: #111827; }}
    .metadata {{ background: #f3f4f6; padding: 1rem; border-radius: 8px; margin-bottom: 1.5rem; }}
    .hero {{ display: block; margin: 2rem auto; max-width: 100%; border-radius: 12px; box-shadow: 0 10px 30px rgba(15, 23, 42, 0.1); }}
    footer {{ margin-top: 3rem; font-size: 0.85rem; text-align: center; color: #6b7280; }}
  </style>
</head>
<body>
  <h1>AI Generated Document</h1>
  <section class=\"metadata\">
    <h2>Summary</h2>
    <p><strong>Prompt:</strong> {prompt}</p>
    <p><strong>Model:</strong> {model}</p>
    <p><strong>Generated at:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
  </section>
  {image_tag}
  <section>
    <h2>Content</h2>
    {paragraphs}
  </section>
  <footer>Generated automatically using {model.split(':')[0].upper()} AI.</footer>
</body>
</html>
"""

    @staticmethod
    def _reset_document(doc: Document) -> None:
        while doc.paragraphs:
            p = doc.paragraphs[0]._element
            p.getparent().remove(p)

    @staticmethod
    def _split_segment(segment: str, index: int) -> tuple[str, str]:
        if "\n" in segment:
            first_line, rest = segment.split("\n", 1)
            title = first_line.strip()[:80] or f"Slide {index}"
            body = rest.strip()[:1500]
        else:
            title = f"Key Insight {index}"
            body = segment[:1500]
        return title, body
